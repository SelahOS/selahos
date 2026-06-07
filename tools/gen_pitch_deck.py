#!/usr/bin/env python3
"""
SelahOS Investor Pitch Deck Generator
14 slides · selahos.io visual identity
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ── Colour palette — mirrors selahos.io CSS variables ─────────────────────────
BG          = RGBColor(0x08, 0x0C, 0x18)   # --bg        #080C18
BG2         = RGBColor(0x0D, 0x12, 0x20)   # --bg2       #0D1220
BG3         = RGBColor(0x13, 0x19, 0x29)   # --bg3       #131929
GOLD        = RGBColor(0xC9, 0x97, 0x4A)   # --gold      #C9974A
GOLD_LT     = RGBColor(0xDE, 0xB9, 0x6E)   # --gold-lt   #DEB96E
PARCH       = RGBColor(0xED, 0xE4, 0xD4)   # --parch     #EDE4D4
PARCH_DIM   = RGBColor(0xB8, 0xAD, 0x9E)   # --parch-dim #B8AD9E
MUTED       = RGBColor(0x6B, 0x63, 0x57)   # --muted     #6B6357
BORDER      = RGBColor(0x1E, 0x25, 0x35)   # --border    #1E2535
TEAL        = RGBColor(0x7D, 0xAD, 0xA3)   # --teal      #7DADA3
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

# ── Slide dimensions (16:9 widescreen) ────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


# ── Low-level helpers ─────────────────────────────────────────────────────────

def solid_fill(shape, color: RGBColor):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    solid_fill(shape, color)
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def add_run(para, text, size_pt, bold=False, color: RGBColor = PARCH, italic=False):
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "DM Sans"
    return run


def add_slide_notes(slide, text: str):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def bg_fill(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def gold_rule_top(slide):
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.055), GOLD)


def gold_rule_bottom(slide):
    add_rect(slide, Inches(0), Inches(7.445), W, Inches(0.055), GOLD)


def section_label(slide, text, color=GOLD):
    tb = add_textbox(slide, Inches(0.55), Inches(0.28), Inches(5), Inches(0.35))
    add_run(tb.text_frame.paragraphs[0], text, 9, bold=True, color=color)


def accent_bar(slide, color=GOLD, width_frac=0.10):
    add_rect(slide, Inches(0.55), Inches(1.08),
             Inches(13.33 * width_frac), Inches(0.04), color)


def footer_line(slide, label="SelahOS  ·  Confidential"):
    tb = add_textbox(slide, Inches(0.55), Inches(7.1), Inches(12.2), Inches(0.3))
    tb.text_frame.word_wrap = False
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    add_run(p, label, 7, color=MUTED)


def heading(slide, text, size=40):
    tb = add_textbox(slide, Inches(0.55), Inches(1.22), Inches(12.2), Inches(0.8))
    add_run(tb.text_frame.paragraphs[0], text, size, bold=True, color=PARCH)
    return tb


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_01_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_fill(slide, BG)
    gold_rule_top(slide)
    gold_rule_bottom(slide)

    # Subtle teal decorative circle
    circle = slide.shapes.add_shape(9, Inches(9.6), Inches(3.8), Inches(4.2), Inches(4.2))
    solid_fill(circle, TEAL)
    circle.line.fill.background()
    sp_pr = circle._element.spPr
    sfe = sp_pr.find(qn('a:solidFill'))
    if sfe is not None:
        srgb = sfe.find(qn('a:srgbClr'))
        if srgb is not None:
            a = etree.SubElement(srgb, qn('a:alpha'))
            a.set('val', '5000')

    # Wordmark
    tb = add_textbox(slide, Inches(0.9), Inches(1.9), Inches(9), Inches(1.4))
    add_run(tb.text_frame.paragraphs[0], "SelahOS", 72, bold=True, color=PARCH)

    # Tagline
    tb2 = add_textbox(slide, Inches(0.9), Inches(3.1), Inches(10), Inches(0.6))
    add_run(tb2.text_frame.paragraphs[0], "Pause. Reflect. Create.", 24,
            italic=True, color=GOLD)

    # Subtitle
    tb3 = add_textbox(slide, Inches(0.9), Inches(3.78), Inches(10), Inches(0.5))
    add_run(tb3.text_frame.paragraphs[0],
            "Human-Centered Computing for the Next Generation", 16, color=PARCH_DIM)

    # Founder line
    tb4 = add_textbox(slide, Inches(0.9), Inches(6.4), Inches(9), Inches(0.5))
    add_run(tb4.text_frame.paragraphs[0],
            "Founder:  Dane-Brandon Noble   ·   Selah Technologies LLC", 11, color=MUTED)

    add_slide_notes(slide,
        "Welcome / opening. Let the tagline land before speaking. "
        "SelahOS is named after the Hebrew musical pause. "
        "Set the tone: mission-driven company, not just a Linux distro.")
    return slide


def slide_02_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_fill(slide, BG)
    accent_bar(slide, GOLD)
    footer_line(slide)
    section_label(slide, "02 / Problem")

    heading(slide, "The Problem")

    tb2 = add_textbox(slide, Inches(0.55), Inches(2.1), Inches(12.2), Inches(0.55))
    tb2.text_frame.word_wrap = True
    add_run(tb2.text_frame.paragraphs[0],
            "Every year, millions of computers are discarded long before their hardware reaches end-of-life.",
            15, italic=True, color=PARCH_DIM)

    cards = [
        ("Rising upgrade costs",         GOLD),
        ("Growing e-waste",               TEAL),
        ("Accessibility challenges",      BG3),
        ("Complex technology ecosystems", BG2),
    ]
    card_w, card_h, card_y, gap = Inches(2.85), Inches(1.55), Inches(2.85), Inches(0.18)
    for i, (label, color) in enumerate(cards):
        x = Inches(0.55) + i * (card_w + gap)
        add_rect(slide, x, card_y, card_w, card_h, color)
        # Gold left border on dark cards
        if color in (BG3, BG2):
            add_rect(slide, x, card_y, Inches(0.05), card_h, GOLD)
        tb_c = add_textbox(slide, x + Inches(0.2), card_y + Inches(0.45),
                           card_w - Inches(0.3), Inches(0.7))
        tb_c.text_frame.word_wrap = True
        p = tb_c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        add_run(p, label, 13, bold=True,
                color=BG if color in (GOLD, TEAL) else PARCH)

    tb3 = add_textbox(slide, Inches(0.55), Inches(5.6), Inches(12.2), Inches(0.5))
    tb3.text_frame.word_wrap = True
    add_run(tb3.text_frame.paragraphs[0],
            "At the same time, millions of users still need affordable, capable computing.",
            14, italic=True, color=MUTED)

    add_slide_notes(slide,
        "Anchor the problem in familiar pain. E-waste and cost resonate with "
        "institutional buyers and ESG-focused investors. "
        "Walk the four cards left to right.")
    return slide


def slide_03_opportunity(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_fill(slide, BG2)
    accent_bar(slide, TEAL, width_frac=0.10)
    footer_line(slide)
    section_label(slide, "03 / Opportunity", TEAL)

    heading(slide, "The Opportunity")

    tb2 = add_textbox(slide, Inches(0.55), Inches(2.15), Inches(12.2), Inches(0.65))
    tb2.text_frame.word_wrap = True
    add_run(tb2.text_frame.paragraphs[0],
            "Most computers are replaced because software outgrows hardware — "
            "not because the hardware has failed.",
            15, italic=True, color=GOLD)

    tb3 = add_textbox(slide, Inches(0.55), Inches(3.05), Inches(6), Inches(0.5))
    add_run(tb3.text_frame.paragraphs[0], "What if we could…", 18, bold=True, color=PARCH)

    bullets = [
        ("Extend hardware life",      GOLD),
        ("Lower costs",               TEAL),
        ("Increase accessibility",    GOLD),
        ("Improve user experience",   TEAL),
    ]
    for i, (text, color) in enumerate(bullets):
        y = Inches(3.65) + i * Inches(0.62)
        add_rect(slide, Inches(0.55), y + Inches(0.12), Inches(0.14), Inches(0.14), color)
        tb_b = add_textbox(slide, Inches(0.85), y, Inches(9), Inches(0.55))
        add_run(tb_b.text_frame.paragraphs[0], text, 17, color=PARCH)

    tb4 = add_textbox(slide, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.45))
    add_run(tb4.text_frame.paragraphs[0],
            "…without requiring new devices?", 15, bold=True, italic=True, color=TEAL)

    add_slide_notes(slide,
        "Pivot from pain to possibility. The rhetorical question invites the audience "
        "into the solution before you name it.")
    return slide


def slide_04_product(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_fill(slide, BG)
    accent_bar(slide, GOLD)
    footer_line(slide)
    section_label(slide, "04 / Solution")

    heading(slide, "Introducing SelahOS")

    tb2 = add_textbox(slide, Inches(0.55), Inches(2.1), Inches(12), Inches(0.5))
    add_run(tb2.text_frame.paragraphs[0],
            "A Human-Centered Computing Platform.", 16, italic=True, color=PARCH_DIM)

    pillars = [
        ("Extend hardware longevity", GOLD),
        ("Improve accessibility",     TEAL),
        ("Empower creators",          GOLD_LT),
        ("Reduce electronic waste",   TEAL),
    ]
    for i, (text, color) in enumerate(pillars):
        row, col = i // 2, i % 2
        bw, bh = Inches(5.9), Inches(1.35)
        bx = Inches(0.55) + col * (bw + Inches(0.28))
        by = Inches(2.8) + row * (bh + Inches(0.22))
        add_rect(slide, bx, by, bw, bh, BG3)
        add_rect(slide, bx, by, Inches(0.055), bh, color)
        tb_p = add_textbox(slide, bx + Inches(0.22), by + Inches(0.38),
                           bw - Inches(0.44), Inches(0.65))
        tb_p.text_frame.word_wrap = True
        add_run(tb_p.text_frame.paragraphs[0], text, 17, bold=True, color=PARCH)

    tb3 = add_textbox(slide, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.45))
    add_run(tb3.text_frame.paragraphs[0],
            "Technology should adapt to people — not the other way around.",
            13, italic=True, color=MUTED)

    add_slide_notes(slide,
        "Name the product, anchor its four core value props. "
        "The philosophy line is the product thesis — say it slowly.")
    return slide


def slide_05_ecosystem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_fill(slide, BG)
    accent_bar(slide, GOLD, width_frac=0.09)
    footer_line(slide)
    section_label(slide, "05 / Ecosystem")

    heading(slide, "The SelahOS Ecosystem")

    tiers = [
        {"name": "Ground", "sub": "Lightweight computing",
         "users": ["Older hardware", "Students", "Seniors", "Community access"],
         "color": TEAL},
        {"name": "Flow", "sub": "Everyday productivity",
         "users": ["Families", "Remote workers", "General users"],
         "color": GOLD},
        {"name": "Studio", "sub": "Creative professionals",
         "users": ["Music", "Design", "Development", "Content creation"],
         "color": BG3},
    ]

    card_w, card_h, card_y, gap = Inches(3.85), Inches(4.3), Inches(2.2), Inches(0.24)
    for i, tier in enumerate(tiers):
        cx = Inches(0.55) + i * (card_w + gap)
        add_rect(slide, cx, card_y, card_w, card_h, BG2)
        add_rect(slide, cx, card_y, Inches(0.055), card_h, tier["color"])

        tb_n = add_textbox(slide, cx + Inches(0.22), card_y + Inches(0.22),
                           card_w - Inches(0.4), Inches(0.55))
        add_run(tb_n.text_frame.paragraphs[0], tier["name"], 22, bold=True,
                color=tier["color"])

        tb_s = add_textbox(slide, cx + Inches(0.22), card_y + Inches(0.78),
                           card_w - Inches(0.4), Inches(0.38))
        add_run(tb_s.text_frame.paragraphs[0], tier["sub"], 11,
                italic=True, color=PARCH_DIM)

        add_rect(slide, cx + Inches(0.22), card_y + Inches(1.22),
                 card_w - Inches(0.44), Inches(0.025), BORDER)

        for j, user in enumerate(tier["users"]):
            uy = card_y + Inches(1.35) + j * Inches(0.62)
            add_rect(slide, cx + Inches(0.25), uy + Inches(0.13),
                     Inches(0.09), Inches(0.09), tier["color"])
            tb_u = add_textbox(slide, cx + Inches(0.45), uy,
                               card_w - Inches(0.65), Inches(0.5))
            add_run(tb_u.text_frame.paragraphs[0], user, 12, color=PARCH_DIM)

    add_slide_notes(slide,
        "Three-tier go-to-market. Ground = underserved; Flow = volume; Studio = revenue. "
        "Present left to right as a natural progression of user sophistication.")
    return slide


def slide_06_selahbridge(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_fill(slide, BG2)
    accent_bar(slide, TEAL)
    footer_line(slide)
    section_label(slide, "06 / Adoption", TEAL)

    heading(slide, "SelahBridge")

    tb2 = add_textbox(slide, Inches(0.55), Inches(2.12), Inches(12.2), Inches(0.55))
    tb2.text_frame.word_wrap = True
    add_run(tb2.text_frame.paragraphs[0],
            "One of the biggest barriers to OS adoption is transition friction. "
            "Users fear losing familiarity.",
            14, italic=True, color=PARCH_DIM)

    boxes = [
        ("Windows\nmacOS", BG3,  "From",  GOLD),
        ("SelahBridge",    TEAL, "Via",   BG),
        ("SelahOS",        GOLD, "To",    BG),
    ]
    bw, bh, by = Inches(3.3), Inches(1.6), Inches(3.0)
    for i, (label, color, caption, text_color) in enumerate(boxes):
        bx = Inches(0.55) + i * (bw + Inches(0.38))
        add_rect(slide, bx, by, bw, bh, color)
        if color == BG3:
            add_rect(slide, bx, by, Inches(0.055), bh, GOLD)
        tb_b = add_textbox(slide, bx + Inches(0.15), by + Inches(0.3),
                           bw - Inches(0.3), bh - Inches(0.4))
        p_b = tb_b.text_frame.paragraphs[0]
        p_b.alignment = PP_ALIGN.CENTER
        add_run(p_b, label, 16, bold=True, color=text_color)
        tb_cap = add_textbox(slide, bx, by + bh + Inches(0.1), bw, Inches(0.3))
        p_cap = tb_cap.text_frame.paragraphs[0]
        p_cap.alignment = PP_ALIGN.CENTER
        add_run(p_cap, caption, 10, color=MUTED)

    outcomes = ["Reduce friction", "Reduce fear", "Increase adoption"]
    for i, text in enumerate(outcomes):
        ty = Inches(5.2) + i * Inches(0.48)
        add_rect(slide, Inches(0.55), ty + Inches(0.12), Inches(0.12), Inches(0.12), TEAL)
        tb_o = add_textbox(slide, Inches(0.82), ty, Inches(8), Inches(0.45))
        add_run(tb_o.text_frame.paragraphs[0], text, 15, bold=True, color=PARCH)

    add_slide_notes(slide,
        "SelahBridge is the key acquisition lever. No other Linux distro ships a "
        "guided migration tool — this is the moat.")
    return slide


def slide_07_neuroadaptive(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_fill(slide, BG)
    accent_bar(slide, GOLD, width_frac=0.08)
    footer_line(slide)
    section_label(slide, "07 / Innovation")

    heading(slide, "Neuroadaptive Computing", size=36)

    tb2 = add_textbox(slide, Inches(0.55), Inches(2.15), Inches(12.2), Inches(0.55))
    tb2.text_frame.word_wrap = True
    add_run(tb2.text_frame.paragraphs[0],
            "Traditional operating systems expect people to adapt. SelahOS inverts that.",
            15, italic=True, color=GOLD)

    features = [
        ("Adaptive workflows",          "UI reconfigures around how you work, not a preset."),
        ("Cognitive-aware experiences", "Reduces overload through contextual simplicity."),
        ("Personalized environments",   "Learns preferences over time — without cloud dependency."),
        ("Creator-focused assistance",  "AI tools embedded where creators already work."),
    ]
    for i, (title, desc) in enumerate(features):
        fy = Inches(3.05) + i * Inches(0.9)
        add_rect(slide, Inches(0.55), fy + Inches(0.18), Inches(0.055), Inches(0.32), TEAL)
        tb_t = add_textbox(slide, Inches(0.75), fy, Inches(5.5), Inches(0.45))
        add_run(tb_t.text_frame.paragraphs[0], title, 14, bold=True, color=PARCH)
        tb_d = add_textbox(slide, Inches(0.75), fy + Inches(0.44), Inches(11.8), Inches(0.38))
        tb_d.text_frame.word_wrap = True
        add_run(tb_d.text_frame.paragraphs[0], desc, 11, color=PARCH_DIM)

    add_slide_notes(slide,
        "Frame as roadmap north star, not vaporware. "
        "Phase 1: heuristic. Phase 2: on-device ML.")
    return slide


def slide_08_impact(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_fill(slide, BG)
    accent_bar(slide, GOLD)
    footer_line(slide)
    section_label(slide, "08 / Impact")

    heading(slide, "Impact")

    pillars = [
        ("Sustainability",  "Reducing e-waste by extending hardware life cycles.",       TEAL),
        ("Accessibility",   "Supporting diverse users — seniors, students, creators.",   GOLD),
        ("Affordability",   "Delivering capable computing without new hardware spend.",  TEAL),
        ("Creativity",      "Giving creators the tools they need on hardware they own.", GOLD),
    ]
    card_w, card_h, card_y, gap = Inches(2.85), Inches(3.5), Inches(2.35), Inches(0.18)
    for i, (title, desc, color) in enumerate(pillars):
        cx = Inches(0.55) + i * (card_w + gap)
        add_rect(slide, cx, card_y, card_w, card_h, BG3)
        add_rect(slide, cx, card_y, Inches(0.055), card_h, color)

        tb_t = add_textbox(slide, cx + Inches(0.2), card_y + Inches(0.25),
                           card_w - Inches(0.35), Inches(0.5))
        add_run(tb_t.text_frame.paragraphs[0], title, 15, bold=True, color=color)

        add_rect(slide, cx + Inches(0.2), card_y + Inches(0.85),
                 card_w - Inches(0.4), Inches(0.025), BORDER)

        tb_d = add_textbox(slide, cx + Inches(0.2), card_y + Inches(1.0),
                           card_w - Inches(0.35), Inches(2.2))
        tb_d.text_frame.word_wrap = True
        add_run(tb_d.text_frame.paragraphs[0], desc, 11, color=PARCH_DIM)

    add_slide_notes(slide,
        "ESG investors, foundations, institutional partners. "
        "Each pillar has a concrete stakeholder behind it.")
    return slide


def slide_09_market(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_fill(slide, BG2)
    accent_bar(slide, TEAL, width_frac=0.09)
    footer_line(slide)
    section_label(slide, "09 / Market", TEAL)

    heading(slide, "Market Opportunity")

    col1 = ["Families", "Schools", "Churches", "Nonprofits",
            "Creators", "Small businesses", "Refurbishment programs"]
    col2 = ["Digital inclusion initiatives", "K-12 and higher education",
            "Accessibility technology", "Emerging / developing markets"]

    tb_h1 = add_textbox(slide, Inches(0.55), Inches(2.2), Inches(5.8), Inches(0.45))
    add_run(tb_h1.text_frame.paragraphs[0], "Target Markets", 13, bold=True, color=GOLD)

    tb_h2 = add_textbox(slide, Inches(7.0), Inches(2.2), Inches(5.8), Inches(0.45))
    add_run(tb_h2.text_frame.paragraphs[0], "Future Expansion", 13, bold=True, color=TEAL)

    for i, item in enumerate(col1):
        y = Inches(2.75) + i * Inches(0.53)
        add_rect(slide, Inches(0.55), y + Inches(0.15), Inches(0.09), Inches(0.09), GOLD)
        tb_i = add_textbox(slide, Inches(0.80), y, Inches(5.5), Inches(0.47))
        add_run(tb_i.text_frame.paragraphs[0], item, 13, color=PARCH_DIM)

    for i, item in enumerate(col2):
        y = Inches(2.75) + i * Inches(0.65)
        add_rect(slide, Inches(7.0), y + Inches(0.15), Inches(0.09), Inches(0.09), TEAL)
        tb_i = add_textbox(slide, Inches(7.25), y, Inches(5.5), Inches(0.55))
        tb_i.text_frame.word_wrap = True
        add_run(tb_i.text_frame.paragraphs[0], item, 13, color=PARCH_DIM)

    add_slide_notes(slide,
        "Refurbishment programs alone = millions of devices annually. "
        "That is the beachhead channel.")
    return slide


def slide_10_why_now(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_fill(slide, BG)
    accent_bar(slide, GOLD)
    footer_line(slide)
    section_label(slide, "10 / Timing")

    heading(slide, "Why Now?")

    tb2 = add_textbox(slide, Inches(0.55), Inches(2.12), Inches(12.2), Inches(0.5))
    tb2.text_frame.word_wrap = True
    add_run(tb2.text_frame.paragraphs[0],
            "Three major trends are converging — creating a window that won't stay open.",
            14, italic=True, color=PARCH_DIM)

    trends = [
        ("01", "Rising hardware costs",
         "Global chip shortages have pushed device prices up 20–35% since 2021.", GOLD),
        ("02", "Growing e-waste concern",
         "53.6M metric tons of e-waste generated in 2024. Only 17% formally recycled.", TEAL),
        ("03", "Demand for accessible technology",
         "Digital equity mandates are creating funded procurement pipelines.", GOLD),
    ]
    for i, (num, title, body, color) in enumerate(trends):
        ty = Inches(2.9) + i * Inches(1.38)
        add_rect(slide, Inches(0.55), ty, Inches(0.55), Inches(1.1), BG3)
        add_rect(slide, Inches(0.55), ty, Inches(0.055), Inches(1.1), color)
        tb_num = add_textbox(slide, Inches(0.55), ty + Inches(0.25), Inches(0.55), Inches(0.55))
        p_num = tb_num.text_frame.paragraphs[0]
        p_num.alignment = PP_ALIGN.CENTER
        add_run(p_num, num, 13, bold=True, color=color)

        tb_t = add_textbox(slide, Inches(1.25), ty + Inches(0.05), Inches(11.3), Inches(0.45))
        add_run(tb_t.text_frame.paragraphs[0], title, 15, bold=True, color=PARCH)

        tb_b = add_textbox(slide, Inches(1.25), ty + Inches(0.5), Inches(11.3), Inches(0.5))
        tb_b.text_frame.word_wrap = True
        add_run(tb_b.text_frame.paragraphs[0], body, 11, color=PARCH_DIM)

    add_slide_notes(slide,
        "Frame each trend as a force already in motion. "
        "SelahOS is positioned at the intersection of all three.")
    return slide


def slide_11_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_fill(slide, BG2)
    accent_bar(slide, GOLD, width_frac=0.09)
    footer_line(slide)
    section_label(slide, "11 / Roadmap")

    heading(slide, "Roadmap")

    phases = [
        {"year": "2026", "theme": "Launch",
         "items": ["Public beta", "Community pilots", "Hardware testing", "SelahBridge"],
         "color": TEAL},
        {"year": "2027", "theme": "Grow",
         "items": ["Education partnerships", "Creator ecosystem", "Accessibility expansion"],
         "color": GOLD},
        {"year": "2028", "theme": "Scale",
         "items": ["Institutional deployments", "Certification programs", "Global scaling"],
         "color": GOLD_LT},
    ]

    add_rect(slide, Inches(0.55), Inches(3.22), Inches(12.2), Inches(0.04), BORDER)

    col_w = Inches(3.95)
    for i, phase in enumerate(phases):
        cx = Inches(0.55) + i * (col_w + Inches(0.22))

        add_rect(slide, cx, Inches(2.15), Inches(1.3), Inches(0.72), phase["color"])
        tb_y = add_textbox(slide, cx + Inches(0.08), Inches(2.2), Inches(1.14), Inches(0.6))
        p_y = tb_y.text_frame.paragraphs[0]
        p_y.alignment = PP_ALIGN.CENTER
        add_run(p_y, phase["year"], 18, bold=True,
                color=BG if phase["color"] in (GOLD, TEAL, GOLD_LT) else PARCH)

        node_cx = cx + Inches(0.65) - Inches(0.12)
        add_rect(slide, node_cx, Inches(3.22) - Inches(0.12),
                 Inches(0.24), Inches(0.24), phase["color"])

        tb_th = add_textbox(slide, cx, Inches(3.55), col_w, Inches(0.45))
        add_run(tb_th.text_frame.paragraphs[0], phase["theme"], 13,
                bold=True, color=phase["color"])

        for j, item in enumerate(phase["items"]):
            iy = Inches(4.1) + j * Inches(0.62)
            add_rect(slide, cx + Inches(0.05), iy + Inches(0.15),
                     Inches(0.09), Inches(0.09), phase["color"])
            tb_i = add_textbox(slide, cx + Inches(0.28), iy, col_w - Inches(0.3), Inches(0.55))
            tb_i.text_frame.word_wrap = True
            add_run(tb_i.text_frame.paragraphs[0], item, 12, color=PARCH_DIM)

    add_slide_notes(slide,
        "2026: validation. 2027: channels. 2028: scale. "
        "Education partnerships in 2027 unlock grant funding.")
    return slide


def slide_12_founder(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_fill(slide, BG)
    accent_bar(slide, GOLD)
    footer_line(slide)
    section_label(slide, "12 / Founder")

    heading(slide, "Dane-Brandon Noble", size=36)

    tb2 = add_textbox(slide, Inches(0.55), Inches(2.08), Inches(8), Inches(0.45))
    add_run(tb2.text_frame.paragraphs[0],
            "Founder & CEO  ·  Selah Technologies LLC", 14, italic=True, color=PARCH_DIM)

    attrs = [
        ("Technology Accessibility Advocate", GOLD),
        ("Creator",                           TEAL),
        ("Community Builder",                 GOLD_LT),
        ("Human-Centered Computing Visionary", TEAL),
    ]
    chip_y, chip_x = Inches(2.72), Inches(0.55)
    for label, color in attrs:
        chip_w = Inches(len(label) * 0.115 + 0.5)
        if chip_x + chip_w > Inches(12.8):
            chip_x  = Inches(0.55)
            chip_y += Inches(0.6)
        add_rect(slide, chip_x, chip_y, chip_w, Inches(0.42), BG3)
        add_rect(slide, chip_x, chip_y, Inches(0.04), Inches(0.42), color)
        tb_c = add_textbox(slide, chip_x + Inches(0.14), chip_y + Inches(0.05),
                           chip_w - Inches(0.22), Inches(0.35))
        add_run(tb_c.text_frame.paragraphs[0], label, 11, bold=True, color=color)
        chip_x += chip_w + Inches(0.15)

    add_rect(slide, Inches(0.55), Inches(4.45), Inches(12.2), Inches(1.4), BG3)
    add_rect(slide, Inches(0.55), Inches(4.45), Inches(0.055), Inches(1.4), GOLD)
    tb_mh = add_textbox(slide, Inches(0.75), Inches(4.58), Inches(11.8), Inches(0.38))
    add_run(tb_mh.text_frame.paragraphs[0], "Mission", 10, bold=True, color=GOLD)
    tb_mb = add_textbox(slide, Inches(0.75), Inches(5.0), Inches(11.8), Inches(0.65))
    tb_mb.text_frame.word_wrap = True
    add_run(tb_mb.text_frame.paragraphs[0],
            "Help people do more with the technology they already own.",
            16, italic=True, color=PARCH)

    add_slide_notes(slide,
        "The founder IS the thesis at this stage. "
        "Share the personal story if the room is small enough.")
    return slide


def slide_13_funding(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_fill(slide, BG2)
    accent_bar(slide, GOLD, width_frac=0.08)
    footer_line(slide)
    section_label(slide, "13 / Use of Funds")

    heading(slide, "What Funding Accelerates", size=36)

    tb2 = add_textbox(slide, Inches(0.55), Inches(2.12), Inches(12.2), Inches(0.45))
    tb2.text_frame.word_wrap = True
    add_run(tb2.text_frame.paragraphs[0],
            "Capital unlocks the speed and scale we cannot achieve organically.",
            13, italic=True, color=PARCH_DIM)

    uses = [
        ("Engineering",              "Core OS, SelahBridge, and Studio tier build-out.", GOLD),
        ("Accessibility Research",   "Neuroadaptive UX research and inclusive design.", TEAL),
        ("Hardware Testing",         "Compatibility lab covering 200+ device configurations.", GOLD),
        ("Documentation",            "Onboarding guides, translated content, video library.", TEAL),
        ("Community Deployments",    "Pilot programs with schools, churches, and nonprofits.", GOLD),
        ("Educational Partnerships", "Curriculum integration and institution licensing.", TEAL),
    ]
    col_w = Inches(5.9)
    for i, (title, desc, color) in enumerate(uses):
        col, row = i % 2, i // 2
        ux = Inches(0.55) + col * (col_w + Inches(0.28))
        uy = Inches(2.75) + row * Inches(1.38)
        add_rect(slide, ux, uy, col_w, Inches(1.15), BG3)
        add_rect(slide, ux, uy, Inches(0.055), Inches(1.15), color)
        tb_t = add_textbox(slide, ux + Inches(0.2), uy + Inches(0.12),
                           col_w - Inches(0.3), Inches(0.42))
        add_run(tb_t.text_frame.paragraphs[0], title, 13, bold=True, color=PARCH)
        tb_d = add_textbox(slide, ux + Inches(0.2), uy + Inches(0.54),
                           col_w - Inches(0.3), Inches(0.48))
        tb_d.text_frame.word_wrap = True
        add_run(tb_d.text_frame.paragraphs[0], desc, 10, color=PARCH_DIM)

    add_slide_notes(slide,
        "Seed range $250K–$1M for pre-revenue open-source OS at this stage. "
        "Each line maps to a measurable roadmap milestone.")
    return slide


def slide_14_cta(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_fill(slide, BG)
    gold_rule_top(slide)
    gold_rule_bottom(slide)

    # Teal decorative circle
    circle = slide.shapes.add_shape(9, Inches(-0.8), Inches(3.5), Inches(4.5), Inches(4.5))
    solid_fill(circle, TEAL)
    circle.line.fill.background()
    sp_pr = circle._element.spPr
    sfe = sp_pr.find(qn('a:solidFill'))
    if sfe is not None:
        srgb = sfe.find(qn('a:srgbClr'))
        if srgb is not None:
            a = etree.SubElement(srgb, qn('a:alpha'))
            a.set('val', '3500')

    section_label(slide, "14 / Join Us")

    tb = add_textbox(slide, Inches(2.0), Inches(1.7), Inches(10), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, "Join Us", 52, bold=True, color=PARCH)

    tb2 = add_textbox(slide, Inches(1.2), Inches(2.8), Inches(11), Inches(0.9))
    tb2.text_frame.word_wrap = True
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    add_run(p2,
            "Building a future where technology lasts longer,\ncosts less, and works better for people.",
            17, italic=True, color=PARCH_DIM)

    tb3 = add_textbox(slide, Inches(1.2), Inches(4.1), Inches(11), Inches(0.85))
    p3 = tb3.text_frame.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    add_run(p3, "SelahOS", 48, bold=True, color=PARCH)

    tb4 = add_textbox(slide, Inches(1.2), Inches(4.95), Inches(11), Inches(0.55))
    p4 = tb4.text_frame.paragraphs[0]
    p4.alignment = PP_ALIGN.CENTER
    add_run(p4, "Pause. Reflect. Create.", 20, italic=True, color=GOLD)

    tb5 = add_textbox(slide, Inches(1.2), Inches(6.15), Inches(11), Inches(0.45))
    p5 = tb5.text_frame.paragraphs[0]
    p5.alignment = PP_ALIGN.CENTER
    add_run(p5, "Dane-Brandon Noble  ·  Selah Technologies LLC", 12, color=MUTED)

    add_slide_notes(slide,
        "Close on conviction, not ask. Let the tagline land, then pause. "
        "The ask comes in conversation. Have a one-pager and data room ready.")
    return slide


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()
    slide_01_cover(prs)
    slide_02_problem(prs)
    slide_03_opportunity(prs)
    slide_04_product(prs)
    slide_05_ecosystem(prs)
    slide_06_selahbridge(prs)
    slide_07_neuroadaptive(prs)
    slide_08_impact(prs)
    slide_09_market(prs)
    slide_10_why_now(prs)
    slide_11_roadmap(prs)
    slide_12_founder(prs)
    slide_13_funding(prs)
    slide_14_cta(prs)

    out_dir  = os.path.expanduser("~/SelahOS-Dev/omolara")
    out_path = os.path.join(out_dir, "SelahOS-Investor-Pitch-Deck.pptx")
    os.makedirs(out_dir, exist_ok=True)
    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()
